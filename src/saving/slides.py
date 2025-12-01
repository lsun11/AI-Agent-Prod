# src/saving/slides.py
import os
import re
import unicodedata
from datetime import datetime
from typing import Any, Iterable
from pptx import Presentation
from pydantic import BaseModel


def _get_items_from_result(result: Any) -> Iterable[Any]:
    """
    Try to extract a list of 'items' from any kind of result model.
    """
    if hasattr(result, "companies") and getattr(result, "companies", None):
        return result.companies

    if hasattr(result, "resources") and getattr(result, "resources", None):
        return result.resources

    if hasattr(result, "items") and getattr(result, "items", None):
        return result.items

    return []


def save_result_slides(query: str, result: Any) -> str:
    """
    Create a .pptx file summarizing the result.
    """
    prs = Presentation()

    # 1) Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = query.strip() or "Research Summary"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # 2) One slide per item
    items = list(_get_items_from_result(result))
    bullet_layout = prs.slide_layouts[1]

    for i, item in enumerate(items, 1):
        slide = prs.slides.add_slide(bullet_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        name = getattr(item, "name", None) or getattr(item, "title", None) or f"Item {i}"
        title.text = name

        def add_bullet(text: str, level: int = 0):
            if not text:
                return
            if not tf.text:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.level = level

        website = getattr(item, "website", None) or getattr(item, "url", None)
        if website:
            add_bullet(f"🌐 Website: {website}")

        if hasattr(item, "pricing_model"):
            pricing_model = getattr(item, "pricing_model", None)
            if pricing_model is not None:
                add_bullet(f"💰 Pricing: {pricing_model}")
        if hasattr(item, "pricing_details"):
            pricing_details = getattr(item, "pricing_details", None)
            if pricing_details:
                add_bullet(f"💰 Details: {pricing_details}")

        if hasattr(item, "is_open_source"):
            is_open = getattr(item, "is_open_source", None)
            if is_open is not None:
                add_bullet(f"📖 Open Source: {is_open}")

        tech_stack = getattr(item, "tech_stack", None)
        if tech_stack:
            add_bullet("🛠 Tech Stack:", level=0)
            for tech in tech_stack[:5]:
                add_bullet(f"- {tech}", level=1)

        competitors = getattr(item, "competitors", None)
        if competitors:
            add_bullet("🤼 Competitors:", level=0)
            for comp in competitors[:5]:
                add_bullet(f"- {comp}", level=1)

        category = getattr(item, "category", None)
        if category:
            add_bullet(f"🧩 Category: {category}")
        tags = getattr(item, "tags", None)
        if tags:
            add_bullet("🏷 Tags:", level=0)
            for tag in tags[:8]:
                add_bullet(f"- {tag}", level=1)

        difficulty = getattr(item, "difficulty", None)
        if difficulty:
            add_bullet(f"📈 Difficulty: {difficulty}")

        description = getattr(item, "description", None) or getattr(item, "summary", None)
        if description and description != "Analysis failed":
            add_bullet("📝 Summary:", level=0)
            add_bullet(description, level=1)

    # 3) Recommendations / Analysis slide
    analysis = getattr(result, "analysis", None) or getattr(result, "summary", None)
    if analysis is not None:
        slide = prs.slides.add_slide(bullet_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        title.text = "Recommendations / Analysis"

        def add_bullet(text: str, level: int = 0):
            if not text:
                return
            if not tf.text:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.level = level

        if isinstance(analysis, str):
            for line in analysis.splitlines():
                stripped = line.strip()
                if not stripped:
                    continue
                add_bullet(stripped)
        else:
            if hasattr(analysis, "model_dump"):
                data = analysis.model_dump()
            elif isinstance(analysis, dict):
                data = analysis
            else:
                add_bullet(str(analysis))
                data = {}

            summary = data.get("summary") or data.get("description")
            if summary:
                add_bullet("Summary:", level=0)
                add_bullet(summary, level=1)

            best_practices = data.get("best_practices") or data.get("tips")
            if best_practices:
                add_bullet("Best Practices:", level=0)
                for bp in best_practices:
                    add_bullet(f"- {bp}", level=1)

            pitfalls = data.get("pitfalls")
            if pitfalls:
                add_bullet("Pitfalls / Risks:", level=0)
                for pf in pitfalls:
                    add_bullet(f"- {pf}", level=1)

            action_plan = data.get("suggested_action_plan") or data.get("action_plan")
            if action_plan:
                add_bullet("Suggested Action Plan:", level=0)
                for step in action_plan:
                    add_bullet(f"- {step}", level=1)

    # 4) Save to disk
    raw_summary = query.strip() or "research"
    normalized = unicodedata.normalize("NFKC", raw_summary)
    safe = re.sub(r'[<>:"/\\|?*\n\r\t]', "_", normalized)
    safe = "_".join(safe.split())[:80]
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{safe}_{timestamp}.pptx"

    folder = "saved_slides"
    os.makedirs(folder, exist_ok=True)

    path = os.path.join(folder, filename)
    prs.save(path)
    return path
