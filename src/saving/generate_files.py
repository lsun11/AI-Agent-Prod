from __future__ import annotations

import os
from pathlib import Path
from typing import Dict

from .layout_llm import DocumentLayout
from .pdf_builder import build_pdf_document
from .docx_builder import build_docx_document


def write_txt(path: str | Path, markdown: str) -> None:
    Path(path).write_text(markdown, encoding="utf-8")


def write_pdf(path: str | Path, layout: DocumentLayout) -> None:
    """
    Use your existing ReportLab logic here, but:
      - use layout.title for the title
      - use layout.report_markdown as the main content
    """
    build_pdf_document(layout.title, layout.report_markdown, str(path))


def write_docx(path: str | Path, layout: DocumentLayout) -> None:
    """
    Use your existing python-docx-based builder.
    Again, feed layout.report_markdown instead of the raw text.
    """
    build_docx_document(str(path), layout.title, layout.report_markdown)


def write_slides(path: str | Path, layout: DocumentLayout) -> None:
    """
    Create a PPTX using python-pptx and layout.slides.
    This can be simpler than your old logic because you already have
    title + bullets nicely separated.
    """
    from pptx import Presentation

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # Title slide
    first = layout.slides[0] if layout.slides else None
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = first.title if first else layout.title
    slide.placeholders[1].text = ""  # subtitle blank or short context

    # Remaining slides
    for slide_data in (layout.slides[1:] if first else layout.slides):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_data.title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(slide_data.bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet

    prs.save(path)


def generate_all_files_for_layout(
    layout: DocumentLayout,
    base_folder: str,
    base_filename: str,
) -> Dict[str, str]:
    """
    Given a DocumentLayout, write txt, pdf, docx, and slides,
    and return their paths.

    NOTE:
    - TXT / PDF / DOCX are always written to project-root "saved_docs/"
    - PPTX is always written to project-root "saved_slides/"
    - `base_folder` is kept for backwards compatibility but ignored
      for the final locations, so that /download/{filename} can
      reliably look in saved_docs/ and saved_slides/.
    """
    docs_dir = Path("saved_docs")
    slides_dir = Path("saved_slides")

    docs_dir.mkdir(parents=True, exist_ok=True)
    slides_dir.mkdir(parents=True, exist_ok=True)

    # Base names (same filename across all formats)
    base_docs = docs_dir / base_filename
    base_slides = slides_dir / base_filename

    txt_path = base_docs.with_suffix(".txt")
    pdf_path = base_docs.with_suffix(".pdf")
    docx_path = base_docs.with_suffix(".docx")
    pptx_path = base_slides.with_suffix(".pptx")

    # Actually write files
    write_txt(txt_path, layout.report_markdown)
    write_pdf(pdf_path, layout)
    write_docx(docx_path, layout)
    write_slides(pptx_path, layout)

    return {
        "txt": str(txt_path),
        "pdf": str(pdf_path),
        "docx": str(docx_path),
        "pptx": str(pptx_path),
    }
