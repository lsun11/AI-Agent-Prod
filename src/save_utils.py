# src/save_utils.py
import os
from datetime import datetime
from typing import Any, Iterable
from langchain_openai import ChatOpenAI
from dotenv import load_dotenv
from pptx import Presentation
from pydantic import BaseModel

from .format_text import to_document

load_dotenv()
## TODO: check if nano is a good model for this highligher.
## Currently it's causing issues.
renderer_llm = ChatOpenAI(model="gpt-4.1-nano", temperature=0)

def _get_items_from_result(result: Any) -> Iterable[Any]:
    """
    Try to extract a list of 'items' from any kind of result model.
    Supports older company-based models and newer generic ones.
    """
    # Original CS shape
    if hasattr(result, "companies") and getattr(result, "companies", None):
        return result.companies

    # Career / other domains might use `resources`
    if hasattr(result, "resources") and getattr(result, "resources", None):
        return result.resources

    # Generic `items` field for some base states
    if hasattr(result, "items") and getattr(result, "items", None):
        return result.items

    # Fallback: empty iterable
    return []


def format_result_text(query: str, result: Any) -> str:
    lines: list[str] = []

    # Generic header
    lines.append(f"📊 Results for: {query}. Click a bubble to open the links")

    items = list(_get_items_from_result(result))
    print("FFFFFFFFFFFFFFFFOOOOOOOOOOOOOOOO", lines, items)
    if not items:
        # If no structured items, just show analysis/summary if available
        lines.append("")
        summary = getattr(result, "analysis", None) or getattr(result, "summary", None)
        if summary:
            lines.append("Summary / Analysis:")
            lines.append("-" * 40)
            lines.append(str(summary))
        else:
            lines.append("(No structured items were returned.)")
        return "\n".join(lines)

    for i, item in enumerate(items, 1):
        # Name / title
        name = getattr(item, "name", None) or getattr(item, "title", None) or f"Item {i}"
        lines.append(f"\n{i}. 🏢 {name}")

        # Website / URL
        website = getattr(item, "website", None) or getattr(item, "url", None)
        if website:
            lines.append(f"   🌐 Website: {website}")

        # Pricing info (CS / SaaS topics)
        if hasattr(item, "pricing_model"):
            pricing_model = getattr(item, "pricing_model", None)
            if pricing_model is not None:
                lines.append(f"   💰 Pricing: {pricing_model}")

        if hasattr(item, "pricing_details"):
            pricing_details = getattr(item, "pricing_details", None)
            if pricing_details:
                lines.append(f"   💰 Pricing Details: {pricing_details}")

        # Open source flag (if applicable)
        if hasattr(item, "is_open_source"):
            is_open = getattr(item, "is_open_source", None)
            if is_open is not None:
                lines.append(f"   📖 Open Source: {is_open}")

        # Tech stack
        tech_stack = getattr(item, "tech_stack", None)
        if tech_stack:
            lines.append(f"   🛠️ Tech Stack: {', '.join(tech_stack[:5])}")

        # Competitors
        competitors = getattr(item, "competitors", None)
        if competitors:
            lines.append(f"   🤼 Competitors: {', '.join(competitors[:5])}")

        # Language support
        language_support = getattr(item, "language_support", None)
        if language_support:
            lines.append(
                f"   💻 Language Support: {', '.join(language_support[:5])}"
            )

        # API availability
        if hasattr(item, "api_available"):
            api_available = getattr(item, "api_available", None)
            if api_available is not None:
                api_status = "✅ Available" if api_available else "❌ Not Available"
                lines.append(f"   🔌 API: {api_status}")

        # Integrations
        integration_caps = getattr(item, "integration_capabilities", None)
        if integration_caps:
            lines.append(
                f"   🔗 Integrations: {', '.join(integration_caps[:4])}"
            )

        target_users = getattr(item, "target_users", None)
        if target_users:
            lines.append(
                f"   🧑‍💻 Target users: {', '.join(target_users[:4])}"
            )

        primary_use_cases = getattr(item, "primary_use_cases", None)
        if primary_use_cases:
            lines.append(
                f"   Primary use cases: {', '.join(primary_use_cases[:4])}"
            )
        # Generic tags / category if present
        category = getattr(item, "category", None)
        if category:
            lines.append(f"   🧩 Category: {category}")

        tags = getattr(item, "tags", None)
        if tags:
            lines.append(f"   🏷️ Tags: {', '.join(tags[:8])}")

        # Difficulty / level (e.g. for learning, career, SE topics)
        difficulty = getattr(item, "difficulty", None)
        if difficulty:
            lines.append(f"   📈 Difficulty: {difficulty}")

        # 🔎 Domain-specific extras (cloud, transportation, etc.)

        # Cloud / infra extras (if present on model)
        if hasattr(item, "free_tier_available"):
            free_tier = getattr(item, "free_tier_available", None)
            if free_tier is not None:
                lines.append(f"   🆓 Free tier available: {free_tier}")

        regions_coverage = getattr(item, "regions_coverage", None)
        if regions_coverage:
            lines.append(f"   🌍 Regions / coverage: {regions_coverage}")

        if hasattr(item, "managed_kubernetes_available"):
            mk = getattr(item, "managed_kubernetes_available", None)
            if mk is not None:
                lines.append(f"   ☸️ Managed Kubernetes: {mk}")

        # Transportation / consumer apps extras
        service_types = getattr(item, "service_types", None)
        if service_types:
            lines.append(f"   🚗 Service types: {', '.join(service_types[:5])}")

        city_coverage = getattr(item, "city_coverage", None)
        if city_coverage:
            lines.append(f"   🌆 City / region coverage: {city_coverage}")

        pricing_model_transport = getattr(item, "pricing_model_transport", None)
        if pricing_model_transport:
            lines.append(f"   💰 Transport pricing: {pricing_model_transport}")

        lines.append("")

        # --- NEW: topic-specific / extra fields ---
        # List of fields we've already displayed and want to skip.
        core_fields = {
            "name", "title", "website", "url",
            "pricing_model", "pricing_details", "is_open_source",
            "tech_stack", "competitors", "api_available",
            "language_support", "integration_capabilities",
            "category", "tags", "difficulty", "description", "summary",
        }

        # Introspect Pydantic model fields if available
        if isinstance(item, BaseModel):
            field_names = item.model_fields.keys()
        else:
            field_names = item.__dict__.keys()

        for field_name in field_names:
            if field_name in core_fields:
                continue
            if field_name.startswith("_"):
                continue

            value = getattr(item, field_name, None)
            if value in (None, [], {}):
                continue

            # Pretty label: service_types -> Service types
            label = field_name.replace("_", " ").capitalize()

            # Simple formatting for lists vs scalars
            if isinstance(value, list):
                rendered = ", ".join(str(v) for v in value[:8])
            else:
                rendered = str(value)

            lines.append(f"   🔧 {label}: {rendered}")

        lines.append("")

    # ... tail part (analysis / to_document / ai_highlight) unchanged ...
    analysis = getattr(result, "analysis", None) or getattr(result, "summary", None)
    if analysis is not None:
        lines.append("**Recommendations / Analysis:**")
        lines.append("-" * 40)
        if isinstance(analysis, str):
            lines.append(to_document(analysis))
        else:
            formatted = to_document(analysis)
            lines.append(formatted)

    analysis_text = "\n".join(lines)
    highlighted = ai_highlight(analysis_text)
    return highlighted




def ai_highlight(text: str) -> str:
    prompt = f"""
You are a formatter. Given the following text:

{text}

Task:
- Do NOT change the wording, punctuation, or sentence order.
- Do NOT add or remove content.
- Wrap the title (e.g "Recommendations") with **double asterisks**.
- Wrap important tools, frameworks, apps, software, and key concepts with **double asterisks**.
- Use Markdown bold: **like this**.
- Return ONLY the modified text.
"""
    response = renderer_llm.invoke(prompt)
    candidate = response.content.strip()

    # --- Safety guard: ensure NO content was removed or changed ---

    def strip_bold(s: str) -> str:
        # Remove all ** so we can compare the underlying text
        return s.replace("**", "")

    base_original = strip_bold(text)
    base_candidate = strip_bold(candidate)

    # If the underlying text changed (content removed/added/reordered),
    # fall back to the original text with NO extra highlighting.
    if base_original != base_candidate:
        # Debug (optional):
        print("ai_highlight: content mismatch, falling back to original")
        return text

    return candidate




def save_result_document_raw(query: str, result_text: str) -> str:
    raw_summary = query.strip() or "research"
    summary_snippet = "_".join(raw_summary.split())[:50]
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{summary_snippet}_{timestamp}.txt"

    folder = "saved_docs"
    os.makedirs(folder, exist_ok=True)

    path = os.path.join(folder, filename)
    with open(path, "w", encoding="utf-8") as f:
        f.write(result_text)

    return path


def save_result_slides(query: str, result: Any) -> str:
    """
    Create a .pptx file:
      - Title slide with the query
      - One slide per item (company/tool/resource)
      - Final slide with recommendations / analysis
    Returns the path to the .pptx file.
    """
    prs = Presentation()

    # ---------- 1) Title slide ----------
    title_slide_layout = prs.slide_layouts[0]  # usually: Title + Subtitle
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = query.strip() or "Research Summary"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # ---------- 2) One slide per item ----------
    items = list(_get_items_from_result(result))

    bullet_layout = prs.slide_layouts[1]  # Title + Content

    for i, item in enumerate(items, 1):
        slide = prs.slides.add_slide(bullet_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        # Title: name / title
        name = getattr(item, "name", None) or getattr(item, "title", None) or f"Item {i}"
        title.text = name

        def add_bullet(text: str, level: int = 0):
            if not text:
                return
            if not tf.text:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.level = level

        # Website
        website = getattr(item, "website", None) or getattr(item, "url", None)
        if website:
            add_bullet(f"🌐 Website: {website}")

        # Pricing
        if hasattr(item, "pricing_model"):
            pricing_model = getattr(item, "pricing_model", None)
            if pricing_model is not None:
                add_bullet(f"💰 Pricing: {pricing_model}")
        if hasattr(item, "pricing_details"):
            pricing_details = getattr(item, "pricing_details", None)
            if pricing_details:
                add_bullet(f"💰 Details: {pricing_details}")

        # Open source
        if hasattr(item, "is_open_source"):
            is_open = getattr(item, "is_open_source", None)
            if is_open is not None:
                add_bullet(f"📖 Open Source: {is_open}")

        # Tech stack
        tech_stack = getattr(item, "tech_stack", None)
        if tech_stack:
            add_bullet("🛠 Tech Stack:", level=0)
            for tech in tech_stack[:5]:
                add_bullet(f"- {tech}", level=1)

        # Competitors
        competitors = getattr(item, "competitors", None)
        if competitors:
            add_bullet("🤼 Competitors:", level=0)
            for comp in competitors[:5]:
                add_bullet(f"- {comp}", level=1)

        # Category / tags
        category = getattr(item, "category", None)
        if category:
            add_bullet(f"🧩 Category: {category}")
        tags = getattr(item, "tags", None)
        if tags:
            add_bullet("🏷 Tags:", level=0)
            for tag in tags[:8]:
                add_bullet(f"- {tag}", level=1)

        # Difficulty
        difficulty = getattr(item, "difficulty", None)
        if difficulty:
            add_bullet(f"📈 Difficulty: {difficulty}")

        # Description / summary
        description = getattr(item, "description", None) or getattr(item, "summary", None)
        if description and description != "Analysis failed":
            add_bullet("📝 Summary:", level=0)
            # You can split into shorter bullets if you want:
            add_bullet(description, level=1)

    # ---------- 3) Recommendations / Analysis slide ----------
    analysis = getattr(result, "analysis", None) or getattr(result, "summary", None)
    if analysis is not None:
        slide = prs.slides.add_slide(bullet_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        title.text = "Recommendations / Analysis"

        def add_bullet(text: str, level: int = 0):
            if not text:
                return
            if not tf.text:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.level = level

        # Case 1: analysis is a plain string
        if isinstance(analysis, str):
            # You can split by lines that look like bullets:
            for line in analysis.splitlines():
                stripped = line.strip()
                if not stripped:
                    continue
                add_bullet(stripped)
        else:
            # Case 2: dict-like / Pydantic
            if hasattr(analysis, "model_dump"):
                data: dict[str, Any] = analysis.model_dump()
            elif isinstance(analysis, dict):
                data = analysis
            else:
                add_bullet(str(analysis))
                data = {}

            summary = data.get("summary") or data.get("description")
            if summary:
                add_bullet("Summary:", level=0)
                add_bullet(summary, level=1)

            best_practices = data.get("best_practices") or data.get("tips")
            if best_practices:
                add_bullet("Best Practices:", level=0)
                for bp in best_practices:
                    add_bullet(f"- {bp}", level=1)

            pitfalls = data.get("pitfalls")
            if pitfalls:
                add_bullet("Pitfalls / Risks:", level=0)
                for pf in pitfalls:
                    add_bullet(f"- {pf}", level=1)

            action_plan = data.get("suggested_action_plan") or data.get("action_plan")
            if action_plan:
                add_bullet("Suggested Action Plan:", level=0)
                for step in action_plan:
                    add_bullet(f"- {step}", level=1)

    # ---------- 4) Save to disk ----------
    raw_summary = query.strip() or "research"
    summary_snippet = "_".join(raw_summary.split())[:50]
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{summary_snippet}_{timestamp}.pptx"

    folder = "saved_slides"
    os.makedirs(folder, exist_ok=True)

    path = os.path.join(folder, filename)
    prs.save(path)
    return path


